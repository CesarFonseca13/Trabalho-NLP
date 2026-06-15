"""
Gera a apresentação (PowerPoint) do trabalho de NLP.
Saída: Apresentacao_NLP.pptx  (13 slides, ~10-15 min, com notas do apresentador).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ----- Paleta (mesma do app) -----
ROXO = RGBColor(0x6C, 0x5C, 0xE7)
CIANO = RGBColor(0x00, 0xC2, 0xC7)
ESCURO = RGBColor(0x1E, 0x1B, 0x2E)
CINZA = RGBColor(0x47, 0x55, 0x69)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
FUNDO_SUAVE = RGBColor(0xF4, 0xF2, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
LARGURA = prs.slide_width
ALTURA = prs.slide_height
BRANCO_LAYOUT = prs.slide_layouts[6]  # slide em branco


def fundo(slide, cor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = cor


def retangulo(slide, x, y, w, h, cor):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = cor
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def caixa_texto(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def paragrafo(tf, texto, tam=18, cor=CINZA, negrito=False, bullet=True,
              space_after=10, primeiro=False, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if primeiro else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = ("•  " + texto) if bullet else texto
    run.font.size = Pt(tam)
    run.font.color.rgb = cor
    run.font.bold = negrito
    run.font.name = "Segoe UI"
    return p


def slide_conteudo(titulo, bullets, notas=""):
    """Slide padrão: barra superior + título + tópicos."""
    s = prs.slides.add_slide(BRANCO_LAYOUT)
    fundo(s, BRANCO)
    retangulo(s, 0, 0, LARGURA, Inches(0.28), ROXO)           # barra de topo
    retangulo(s, 0, Inches(0.28), Inches(2.4), Inches(0.06), CIANO)

    tf_t = caixa_texto(s, Inches(0.7), Inches(0.55), Inches(12), Inches(1.0))
    paragrafo(tf_t, titulo, tam=32, cor=ESCURO, negrito=True, bullet=False, primeiro=True)

    tf = caixa_texto(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(5.2))
    for i, b in enumerate(bullets):
        if isinstance(b, tuple):
            texto, tam, negrito = b
            paragrafo(tf, texto, tam=tam, cor=ESCURO if negrito else CINZA,
                      negrito=negrito, primeiro=(i == 0))
        else:
            paragrafo(tf, b, tam=20, primeiro=(i == 0))
    if notas:
        s.notes_slide.notes_text_frame.text = notas
    return s


def slide_imagem(titulo, caminho_img, legenda, notas=""):
    s = prs.slides.add_slide(BRANCO_LAYOUT)
    fundo(s, BRANCO)
    retangulo(s, 0, 0, LARGURA, Inches(0.28), ROXO)
    retangulo(s, 0, Inches(0.28), Inches(2.4), Inches(0.06), CIANO)

    tf_t = caixa_texto(s, Inches(0.7), Inches(0.55), Inches(12), Inches(1.0))
    paragrafo(tf_t, titulo, tam=32, cor=ESCURO, negrito=True, bullet=False, primeiro=True)

    if os.path.exists(caminho_img):
        # centraliza a imagem mantendo proporção dentro de uma área
        from PIL import Image
        try:
            iw, ih = Image.open(caminho_img).size
            ratio = iw / ih
        except Exception:
            ratio = 1.6
        max_h = Inches(4.7)
        h = max_h
        w = int(h * ratio)
        if w > Inches(11):
            w = Inches(11)
            h = int(w / ratio)
        x = int((LARGURA - w) / 2)
        s.shapes.add_picture(caminho_img, x, Inches(1.7), width=w, height=h)

    tf = caixa_texto(s, Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.7))
    paragrafo(tf, legenda, tam=16, cor=ROXO, negrito=True, bullet=False,
              primeiro=True, align=PP_ALIGN.CENTER)
    if notas:
        s.notes_slide.notes_text_frame.text = notas
    return s


# ============================================================
# SLIDE 1 — CAPA
# ============================================================
s = prs.slides.add_slide(BRANCO_LAYOUT)
fundo(s, ROXO)
retangulo(s, 0, Inches(5.0), LARGURA, Inches(0.12), CIANO)
tf = caixa_texto(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.0))
paragrafo(tf, "Busca Semântica em Textos", tam=46, cor=BRANCO, negrito=True,
          bullet=False, primeiro=True)
paragrafo(tf, "Recuperação inteligente de informação com NLP e embeddings vetoriais",
          tam=22, cor=RGBColor(0xEA, 0xE7, 0xFF), bullet=False, space_after=4)
tf2 = caixa_texto(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.5))
paragrafo(tf2, "Processamento de Linguagem Natural  •  IESB",
          tam=18, cor=BRANCO, bullet=False, primeiro=True)
paragrafo(tf2, "Cesar Eustaquio da Fonseca Filho  —  2026",
          tam=18, cor=RGBColor(0xEA, 0xE7, 0xFF), bullet=False)
s.notes_slide.notes_text_frame.text = (
    "Apresente-se em uma frase. Diga que o trabalho mostra como a IA pode buscar "
    "informação por SIGNIFICADO, e não apenas por palavras exatas. (~30s)"
)

# ============================================================
# SLIDE 2 — O PROBLEMA
# ============================================================
slide_conteudo(
    "O problema",
    [
        ("A busca tradicional procura PALAVRAS EXATAS.", 22, True),
        "Se você busca \"celular travando\" e o texto diz \"aparelho congela\", a busca por palavra-chave não encontra.",
        "Em grandes volumes de texto, isso faz perder informação relevante.",
        ("Precisamos buscar por SIGNIFICADO, não por palavra exata.", 22, True),
    ],
    notas=(
        "Explique a dor: busca por palavra-chave é literal. Dê o exemplo "
        "'celular travando' vs 'aparelho congela' — mesmo sentido, palavras diferentes. "
        "É esse problema que o trabalho resolve. (~1 min)"
    ),
)

# ============================================================
# SLIDE 3 — OBJETIVO
# ============================================================
slide_conteudo(
    "Objetivo do trabalho",
    [
        ("Desenvolver um sistema de NLP capaz de:", 22, True),
        "ANALISAR um grande volume de textos;",
        "ORGANIZAR esses textos por similaridade de conteúdo;",
        "RECUPERAR informação de forma semântica (por significado).",
        ("Demonstrando como embeddings vetoriais melhoram a busca em texto.", 20, False),
    ],
    notas=(
        "Leia os três verbos: analisar, organizar, recuperar. Diga que o sistema usa "
        "técnicas modernas de IA (embeddings) para isso. (~45s)"
    ),
)

# ============================================================
# SLIDE 4 — O QUE SÃO EMBEDDINGS
# ============================================================
slide_conteudo(
    "A ideia central: embeddings",
    [
        ("Embedding = transformar um texto em um vetor de números.", 22, True),
        "Cada frase vira um ponto em um espaço de muitas dimensões (aqui, 384).",
        "Textos com SIGNIFICADO parecido ficam PRÓXIMOS nesse espaço.",
        "\"chegou quebrado\" e \"veio com defeito\" → vetores próximos.",
        ("Assim, medir distância entre vetores = medir semelhança de sentido.", 20, True),
    ],
    notas=(
        "Essa é a parte conceitual mais importante. Use a analogia: cada texto vira um "
        "ponto num mapa; textos parecidos ficam perto. Não precisa entrar em matemática. (~1,5 min)"
    ),
)

# ============================================================
# SLIDE 5 — OS DADOS
# ============================================================
slide_conteudo(
    "Os dados",
    [
        ("Base B2W-Reviews01 — avaliações reais de e-commerce em português.", 22, True),
        "≈ 130 mil avaliações de produtos (Americanas/B2W).",
        "Colunas usadas: texto da avaliação, categoria, nota (1–5) e título.",
        "Texto real, escrito por clientes — com gírias, erros e abreviações.",
    ],
    notas=(
        "Diga que escolheu uma base REAL e em português, para o trabalho ter valor "
        "prático. Mencione o tamanho (~130 mil). (~45s)"
    ),
)

# ============================================================
# SLIDE 6 — METODOLOGIA / PIPELINE
# ============================================================
slide_conteudo(
    "Metodologia (passo a passo)",
    [
        ("1.  Limpeza do texto — minúsculas, sem pontuação, remoção de stopwords.", 20, False),
        ("2.  Análise exploratória — palavras e expressões mais frequentes.", 20, False),
        ("3.  Embeddings — cada avaliação vira um vetor (modelo multilíngue).", 20, False),
        ("4.  Organização — agrupamento (K-Means) por similaridade.", 20, False),
        ("5.  Busca semântica — consulta → vetor → similaridade do cosseno.", 20, True),
    ],
    notas=(
        "Passe rápido pelas 5 etapas, sem aprofundar. Diga que vai mostrar os resultados "
        "de cada uma nos próximos slides. (~1 min)"
    ),
)

# ============================================================
# SLIDE 7 — EDA (imagem: palavras frequentes)
# ============================================================
slide_imagem(
    "Análise: palavras mais frequentes",
    "img_frequencia.png",
    "Termos como 'produto', 'entrega', 'recomendo' dominam as avaliações.",
    notas=(
        "Mostre que, antes da parte semântica, dá para entender a base só pela frequência "
        "das palavras. Comente 1 ou 2 termos do gráfico. (~45s)"
    ),
)

# ============================================================
# SLIDE 8 — COMO A BUSCA FUNCIONA
# ============================================================
slide_conteudo(
    "Como a busca semântica funciona",
    [
        ("1.  O usuário digita uma consulta em linguagem natural.", 20, False),
        ("2.  A consulta é transformada no mesmo tipo de vetor das avaliações.", 20, False),
        ("3.  Calculamos a similaridade do cosseno entre a consulta e cada texto.", 20, False),
        ("4.  Retornamos as avaliações mais próximas (maior similaridade).", 20, False),
        ("Similaridade do cosseno = o quão alinhados dois vetores estão (0 a 1).", 18, True),
    ],
    notas=(
        "Explique o cosseno de forma simples: quanto mais alinhados os vetores, mais "
        "parecido o sentido; valor perto de 1 = muito parecido. (~1 min)"
    ),
)

# ============================================================
# SLIDE 9 — RESULTADO DA BUSCA (o ponto alto)
# ============================================================
slide_conteudo(
    "Resultado: busca por significado",
    [
        ("Consulta: \"o produto chegou quebrado e com defeito\"", 22, True),
        "→ \"O produto veio quebrado em várias partes\"   (83%)",
        "→ \"O produto veio furado, não dá pra aproveitar\"   (78%)",
        ("Repare: encontrou \"furado\" e \"defeito\" sem a consulta usar essas palavras.", 19, True),
        ("Consulta: \"a bateria descarrega muito rápido\"  →  celulares c/ bateria fraca (87%)", 19, False),
    ],
    notas=(
        "ESTE é o slide mais importante — é a prova de que funciona. Destaque que a busca "
        "achou 'furado' e 'bateria fraca' por SENTIDO, não por palavra. Pode abrir o app "
        "ao vivo aqui se quiser. (~1,5 min)"
    ),
)

# ============================================================
# SLIDE 10 — ORGANIZAÇÃO (imagem: clusters)
# ============================================================
slide_imagem(
    "Organização: agrupamento semântico",
    "img_clusters.png",
    "O K-Means separa as avaliações em grupos de conteúdo parecido, automaticamente.",
    notas=(
        "Mostre que, além de buscar, o sistema ORGANIZA: cada cor é um grupo de avaliações "
        "com assunto parecido, descoberto sem rótulos. (~45s)"
    ),
)

# ============================================================
# SLIDE 11 — O SISTEMA ONLINE
# ============================================================
s11 = slide_conteudo(
    "O sistema online (Streamlit)",
    [
        ("Tudo isso virou um app web simples de usar.", 22, True),
        "O usuário digita a busca e vê os resultados em cards.",
        "Cada resultado mostra: categoria, nota, texto e % de similaridade.",
        "Publicado gratuitamente na nuvem (Streamlit Community Cloud).",
    ],
    notas=(
        "Aqui você abre o app ao vivo (ou mostra um print). Faça UMA busca na frente da "
        "turma — é o momento de maior impacto. O link está no slide. (~1,5 min)"
    ),
)

# Caixa de destaque com o link clicável do app
LINK_APP = "https://trabalho-nlp-lw2nakkbxysdjuxkaipr3x.streamlit.app/"
cx = retangulo(s11, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.95), ROXO)
cx.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
cx.text_frame.word_wrap = True
p_link = cx.text_frame.paragraphs[0]
p_link.alignment = PP_ALIGN.CENTER
r1 = p_link.add_run()
r1.text = "🔗  Acesse o app:  "
r1.font.size = Pt(18); r1.font.bold = True; r1.font.color.rgb = BRANCO; r1.font.name = "Segoe UI"
r2 = p_link.add_run()
r2.text = LINK_APP
r2.font.size = Pt(18); r2.font.color.rgb = RGBColor(0xEA, 0xE7, 0xFF); r2.font.name = "Segoe UI"
r2.hyperlink.address = LINK_APP

# ============================================================
# SLIDE 12 — POTENCIAL DE APLICAÇÃO
# ============================================================
slide_conteudo(
    "Potencial: a memória da empresa",
    [
        ("E se a empresa pudesse \"lembrar\" de tudo o que já produziu?", 22, True),
        "Alimente o sistema com qualquer texto: documentos, projetos, pesquisas científicas, atas e conversas com investidores.",
        "Tudo vira vetor e fica recuperável por significado — uma verdadeira memória semântica da organização.",
        ("Exemplos de uso:", 20, True),
        "Resgatar um achado de uma pesquisa científica interna para um projeto em andamento.",
        "Relembrar o que foi tratado em conversas com investidores.",
        "Encontrar projetos antigos parecidos com o desafio de agora.",
    ],
    notas=(
        "Apresente como a VISÃO de futuro do trabalho. A mesma técnica vira um 'banco de "
        "memórias' da empresa: em vez de avaliações, indexe documentos, pesquisas e atas. "
        "Aí qualquer pessoa 'lembra' de algo buscando por significado — ex.: puxar um dado "
        "de uma pesquisa antiga para um projeto atual. É a base do que o mercado chama de "
        "RAG (busca semântica + IA). (~1 min)"
    ),
)

# ============================================================
# SLIDE 13 — CONCLUSÃO
# ============================================================
slide_conteudo(
    "Conclusão",
    [
        ("Demonstramos uma busca que entende o SIGNIFICADO do texto.", 22, True),
        "Embeddings + similaridade do cosseno recuperam informação relevante.",
        "É a base de tecnologias atuais: buscadores semânticos, assistentes de IA e bancos de dados vetoriais (RAG).",
        ("NLP ajuda a interpretar e recuperar conhecimento em grandes volumes de texto.", 20, True),
    ],
    notas=(
        "Feche conectando com o mundo real: é a mesma ideia por trás do ChatGPT com busca, "
        "do Google semântico, etc. Reforce que o objetivo foi cumprido. (~1 min)"
    ),
)

# ============================================================
# SLIDE 13 — OBRIGADO
# ============================================================
s = prs.slides.add_slide(BRANCO_LAYOUT)
fundo(s, ROXO)
retangulo(s, 0, Inches(4.7), LARGURA, Inches(0.12), CIANO)
tf = caixa_texto(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.0), anchor=MSO_ANCHOR.MIDDLE)
paragrafo(tf, "Obrigado!", tam=48, cor=BRANCO, negrito=True, bullet=False, primeiro=True)
s.notes_slide.notes_text_frame.text = "Agradeça e abra para perguntas. (~30s)"

prs.save("Apresentacao_NLP.pptx")
print("Apresentação gerada: Apresentacao_NLP.pptx  |  slides:", len(prs.slides._sldIdLst))
